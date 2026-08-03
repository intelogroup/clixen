"""
Document creation tools — generate DOCX, PDF, XLSX, and PPTX from AI data.

Supports: Markdown → DOCX/PDF, JSON → DOCX/XLSX, Markdown → PPTX
Optimized for Gemma4 flat-argument tool calling (no deeply nested structures).
"""

from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _read_markdown(path: str) -> str:
    """Read markdown file content."""
    return Path(path).expanduser().read_text(encoding="utf-8")


def _write_file(path: str, content: str) -> str:
    """Write content to file, return path."""
    p = Path(path).expanduser()
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content, encoding="utf-8")
    return str(p)


# ---------------------------------------------------------------------------
# 1. Markdown → DOCX
# ---------------------------------------------------------------------------

def markdown_to_docx(md_path: str, output_path: str, template: str = "") -> str:
    """
    Convert a Markdown file to a Word (.docx) document.

    Args:
        md_path: Path to input Markdown file
        output_path: Path for output .docx file
        template: Optional .docx template path for styling

    Returns:
        Path to created .docx file
    """
    import markdown as _md
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

    md_text = _read_markdown(md_path)
    html = _md.markdown(md_text, extensions=["extra", "codehilite", "tables"])
    # Parse HTML into document elements
    from xml.etree import ElementTree as ET
    root = ET.fromstring(f"<root>{html}</root>")

    doc = Document()
    if template and Path(template).exists():
        doc = Document(template)

    style_map = {
        "h1": lambda d, t: _add_heading(d, t, 0),
        "h2": lambda d, t: _add_heading(d, t, 1),
        "h3": lambda d, t: _add_heading(d, t, 2),
        "p": lambda d, t: d.add_paragraph(t),
        "pre": lambda d, t: _add_code_block(d, t),
        "table": lambda d, t: _add_table(d, t),
    }

    for elem in root:
        tag = elem.tag
        text = _elem_text(elem)
        if not text.strip():
            continue
        handler = style_map.get(tag)
        if handler:
            handler(doc, text)
        else:
            doc.add_paragraph(text)

    output_p = Path(output_path).expanduser()
    output_p.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_p))
    return str(output_p)


def _add_heading(doc, text: str, level: int):
    doc.add_heading(text, level=level)


def _add_code_block(doc, text: str):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.font.name = "Courier New"
    run.font.size = Pt(9)


def _add_table(doc, html_table: str):
    from xml.etree import ElementTree as ET
    t = ET.fromstring(f"<root>{html_table}</root>")
    rows = t.findall(".//tr")
    if not rows:
        return
    cols = len(rows[0].findall("td") or rows[0].findall("th"))
    table = doc.add_table(rows=len(rows), cols=cols)
    for i, row in enumerate(rows):
        cells = row.findall("td") or row.findall("th")
        for j, cell in enumerate(cells[:cols]):
            table.rows[i].cells[j].text = _elem_text(cell)


def _elem_text(elem) -> str:
    from xml.etree import ElementTree as ET
    return "".join(elem.itertext()).strip()


# ---------------------------------------------------------------------------
# 2. Markdown → PDF (via ReportLab)
# ---------------------------------------------------------------------------

def markdown_to_pdf(md_path: str, output_path: str, theme: str = "light") -> str:
    """
    Convert a Markdown file to PDF using ReportLab.

    Args:
        md_path: Path to input Markdown file
        output_path: Path for output .pdf file
        theme: Visual theme — 'light' or 'dark'

    Returns:
        Path to created .pdf file
    """
    import markdown as _md
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Preformatted
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors

    md_text = _read_markdown(md_path)
    html = _md.markdown(md_text, extensions=["extra", "tables"])
    from xml.etree import ElementTree as ET
    root = ET.fromstring(f"<root>{html}</root>")

    doc = SimpleDocTemplate(output_path)
    styles = getSampleStyleSheet()
    if theme == "dark":
        styles.add(ParagraphStyle(name="DarkBody", textColor=colors.white, backColor=colors.black))
    elements = []

    for elem in root:
        tag = elem.tag
        text = _elem_text(elem)
        if not text.strip():
            continue
        if tag in ("h1", "h2", "h3"):
            level = int(tag[1])
            elements.append(Paragraph(text, styles[f"Heading{level}"]))
        elif tag == "pre":
            elements.append(Preformatted(text, styles["Code"]))
        elif tag == "p":
            elements.append(Paragraph(text, styles["BodyText"]))
        elif tag in ("ul", "ol"):
            for i, li in enumerate(elem.findall("li"), start=1):
                li_text = _elem_text(li)
                if not li_text.strip():
                    continue
                prefix = f"{i}. " if tag == "ol" else "• "
                elements.append(Paragraph(prefix + li_text, styles["BodyText"]))
        elements.append(Spacer(1, 0.2 * inch))

    doc.build(elements)
    return output_path


def create_pdf(content: str, output_path: str, theme: str = "light") -> str:
    """
    Create a PDF directly from Markdown-formatted text — combines write + convert
    into one call so the agent doesn't need a separate write_file round trip.

    Args:
        content: Markdown-formatted text (headings, bullet lists, paragraphs)
        output_path: Output .pdf file path
        theme: Visual theme — 'light' or 'dark'

    Returns:
        Path to created .pdf file
    """
    out = Path(output_path).expanduser()
    md_path = out.with_suffix(".md")
    _write_file(str(md_path), content)
    return markdown_to_pdf(str(md_path), str(out), theme=theme)


# ---------------------------------------------------------------------------
# 3. JSON → DOCX (template-style with json2docx patterns)
# ---------------------------------------------------------------------------

def json_to_docx(json_path: str, output_path: str, template: str = "") -> str:
    """
    Convert a JSON data file to a Word (.docx) document.

    JSON structure:
    {
        "title": "Report Title",
        "sections": [{"heading": "...", "content": "..."}]
    }

    Args:
        json_path: Path to input JSON file
        output_path: Path for output .docx file
        template: Optional .docx template path

    Returns:
        Path to created .docx file
    """
    from docx import Document

    raw = Path(json_path).expanduser().read_text(encoding="utf-8")
    data = json.loads(raw)

    doc = Document()
    if template and Path(template).exists():
        doc = Document(template)

    # Title
    title = data.get("title", Path(json_path).expanduser().stem)
    doc.add_heading(title, 0)

    # Sections
    for section in data.get("sections", []):
        heading = section.get("heading", "")
        content = section.get("content", "")
        if heading:
            doc.add_heading(heading, level=1)
        if content:
            doc.add_paragraph(content)

    # Tables (if present)
    for table_data in data.get("tables", []):
        doc.add_heading(table_data.get("name", "Table"), level=2)
        rows = table_data.get("rows", [])
        if rows:
            cols = len(rows[0])
            table = doc.add_table(rows=len(rows), cols=cols)
            for i, row in enumerate(rows):
                for j, cell in enumerate(row[:cols]):
                    table.rows[i].cells[j].text = str(cell)

    output_p = Path(output_path).expanduser()
    output_p.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_p))
    return str(output_p)


# ---------------------------------------------------------------------------
# 4. Data → XLSX
# ---------------------------------------------------------------------------

def data_to_xlsx(data_path: str, output_path: str, sheet_name: str = "Sheet1") -> str:
    """
    Create an Excel (.xlsx) file from JSON or CSV data.

    JSON structure (list of dicts or dict of lists):
    [
        {"Name": "Alice", "Age": 30},
        {"Name": "Bob", "Age": 25}
    ]

    Args:
        data_path: Path to JSON or CSV input file
        output_path: Path for output .xlsx file
        sheet_name: Name of the worksheet

    Returns:
        Path to created .xlsx file
    """
    import openpyxl

    raw = Path(data_path).expanduser().read_text(encoding="utf-8")
    ext = Path(data_path).expanduser().suffix.lower()

    if ext == ".csv":
        import csv
        reader = csv.DictReader(raw.splitlines())
        data = list(reader)
    else:
        data = json.loads(raw)

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name

    if isinstance(data, list) and data:
        headers = list(data[0].keys())
        ws.append(headers)
        for row in data:
            ws.append([row.get(h, "") for h in headers])
    elif isinstance(data, dict):
        for key, value in data.items():
            ws.append([key, str(value)])

    output_p = Path(output_path).expanduser()
    output_p.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(output_p))
    return str(output_p)


# ---------------------------------------------------------------------------
# 5. Markdown → PPTX
# ---------------------------------------------------------------------------

def markdown_to_pptx(md_path: str, output_path: str, template: str = "") -> str:
    """
    Convert a Markdown file to a PowerPoint (.pptx) presentation.

    Markdown structure:
    - # Heading 1 starts a new slide.
    - ## Heading 2 becomes slide title.
    - Bullet lists become slide content.
    - Code blocks are rendered as images (via Pygments).

    Args:
        md_path: Path to input Markdown file.
        output_path: Path for output .pptx file.
        template: Optional .pptx template path.

    Returns:
        Path to created .pptx file.
    """
    try:
        import pptx
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"
    
    import markdown as _md
    from pptx.util import Inches, Pt

    md_text = _read_markdown(md_path)

    prs = pptx.Presentation()
    if template and Path(template).exists():
        prs = pptx.Presentation(template)

    current_slide = None
    current_content = []

    def _finalize_slide():
        nonlocal current_slide, current_content
        if current_slide and current_content:
            textbox = current_slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(8), Inches(4)
            )
            textbox.text = "\n".join(current_content)
        current_content = []

    lines = md_text.splitlines()
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            _finalize_slide()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = line[2:].strip()
            current_slide = slide
        elif line.startswith("## "):
            if current_slide:
                current_slide.shapes.title.text = line[3:].strip()
        elif line.startswith("- ") or line.startswith("* "):
            if current_slide:
                current_content.append(line[2:].strip())
        elif current_slide:
            current_content.append(line)

    _finalize_slide()
    output_p = Path(output_path).expanduser()
    output_p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_p))
    return str(output_p)


# ---------------------------------------------------------------------------
# 6. Text/HTML file creation (simple, no external deps)
# ---------------------------------------------------------------------------

def text_to_file(content: str, output_path: str, encoding: str = "utf-8") -> str:
    """
    Create a plain text (.txt) or HTML (.html) file from string content.

    Args:
        content: Text/HTML content to write
        output_path: Path for output file (.txt or .html)
        encoding: File encoding (default utf-8)

    Returns:
        Path to created file
    """
    p = Path(output_path).expanduser()
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content, encoding=encoding)
    return str(p)


def html_to_file(html_content: str, output_path: str) -> str:
    """
    Create an HTML file from string content.
    Convenience wrapper around text_to_file for clarity.

    Args:
        html_content: HTML markup as a string
        output_path: Path for output .html file

    Returns:
        Path to created .html file
    """
    if not output_path.lower().endswith(".html"):
        output_path = output_path + ".html"
    return text_to_file(html_content, output_path)


# ---------------------------------------------------------------------------
# Tool schemas (Gemma4-optimized — flat arguments)
# ---------------------------------------------------------------------------

MARKDOWN_TO_DOCX_SCHEMA = {
    "type": "function",
    "function": {
        "name": "markdown_to_docx",
        "description": (
            "Create a Word (.docx) document from a Markdown file. "
            "Use when the user asks to create a document, report, letter, or summary from markdown text. "
            "Supports headings, lists, tables, and code blocks."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "md_path": {"type": "string", "description": "Path to input Markdown file"},
                "output_path": {"type": "string", "description": "Output .docx file path"},
                "template": {"type": "string", "default": "", "description": "Optional .docx template for styling"},
            },
            "required": ["md_path", "output_path"],
        },
    },
}

MARKDOWN_TO_PDF_SCHEMA = {
    "type": "function",
    "function": {
        "name": "markdown_to_pdf",
        "description": (
            "Create a PDF document from a Markdown file using ReportLab. "
            "Use when the user asks to generate a PDF report, documentation, or printable document from markdown."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "md_path": {"type": "string", "description": "Path to input Markdown file"},
                "output_path": {"type": "string", "description": "Output .pdf file path"},
                "theme": {"type": "string", "default": "light", "description": "Visual theme: 'light' or 'dark'"},
            },
            "required": ["md_path", "output_path"],
        },
    },
}

CREATE_PDF_SCHEMA = {
    "type": "function",
    "function": {
        "name": "create_pdf",
        "description": (
            "Create a PDF directly from Markdown-formatted text in ONE call — no need to "
            "write a markdown file first. Use this instead of write_file + markdown_to_pdf "
            "whenever the content is being generated fresh (not read from an existing file)."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "content": {"type": "string", "description": "Markdown-formatted text: headings (#), bullet lists (*), paragraphs"},
                "output_path": {"type": "string", "description": "Output .pdf file path"},
                "theme": {"type": "string", "default": "light", "description": "Visual theme: 'light' or 'dark'"},
            },
            "required": ["content", "output_path"],
        },
    },
}

JSON_TO_DOCX_SCHEMA = {
    "type": "function",
    "function": {
        "name": "json_to_docx",
        "description": (
            "Create a Word (.docx) document from a JSON data file. "
            "JSON should have 'title' and 'sections' (list of {heading, content}) or 'tables'. "
            "Use for structured reports from AI-generated JSON."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "json_path": {"type": "string", "description": "Path to input JSON file"},
                "output_path": {"type": "string", "description": "Output .docx file path"},
                "template": {"type": "string", "default": "", "description": "Optional .docx template path"},
            },
            "required": ["json_path", "output_path"],
        },
    },
}

DATA_TO_XLSX_SCHEMA = {
    "type": "function",
    "function": {
        "name": "data_to_xlsx",
        "description": (
            "Create an Excel (.xlsx) file from JSON or CSV data. "
            "JSON can be a list of objects (rows) or a dict of key-value pairs. "
            "Use when the user asks to export data to Excel from AI output."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "data_path": {"type": "string", "description": "Path to JSON or CSV input file"},
                "output_path": {"type": "string", "description": "Output .xlsx file path"},
                "sheet_name": {"type": "string", "default": "Sheet1", "description": "Worksheet name"},
            },
            "required": ["data_path", "output_path"],
        },
    },
}

MARKDOWN_TO_PPTX_SCHEMA = {
    "type": "function",
    "function": {
        "name": "markdown_to_pptx",
        "description": (
            "Create a PowerPoint (.pptx) presentation from a Markdown file. "
            "Use # for slide titles, ## for slide headings, - or * for bullet points. "
            "Use when the user asks to create a presentation or slides from markdown content."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "md_path": {"type": "string", "description": "Path to input Markdown file"},
                "output_path": {"type": "string", "description": "Output .pptx file path"},
                "template": {"type": "string", "default": "", "description": "Optional .pptx template path"},
            },
            "required": ["md_path", "output_path"],
        },
    },
}

TEXT_TO_FILE_SCHEMA = {
    "type": "function",
    "function": {
        "name": "text_to_file",
        "description": (
            "Create a plain text (.txt) or HTML (.html) file from string content. "
            "Use when the user asks to save text, notes, logs, or HTML to a file. "
            "Supports any extension: .txt, .html, .csv, .json, .md, etc."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "content": {"type": "string", "description": "Text content to write to the file"},
                "output_path": {"type": "string", "description": "Output file path (e.g. /tmp/notes.txt)"},
                "encoding": {"type": "string", "default": "utf-8", "description": "File encoding"},
            },
            "required": ["content", "output_path"],
        },
    },
}

HTML_TO_FILE_SCHEMA = {
    "type": "function",
    "function": {
        "name": "html_to_file",
        "description": (
            "Create an HTML file from string content. "
            "Use when the user asks to create a webpage, HTML report, or save HTML markup to a file."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "html_content": {"type": "string", "description": "HTML markup as a string"},
                "output_path": {"type": "string", "description": "Output .html file path"},
            },
            "required": ["html_content", "output_path"],
        },
    },
}


# ---------------------------------------------------------------------------
# Executors (for registry.py)
# ---------------------------------------------------------------------------

def markdown_to_docx_executor(args: dict) -> str:
    return markdown_to_docx(
        md_path=args["md_path"],
        output_path=args["output_path"],
        template=args.get("template", ""),
    )


def markdown_to_pdf_executor(args: dict) -> str:
    return markdown_to_pdf(
        md_path=args["md_path"],
        output_path=args["output_path"],
        theme=args.get("theme", "light"),
    )


def create_pdf_executor(args: dict) -> str:
    return create_pdf(
        content=args["content"],
        output_path=args["output_path"],
        theme=args.get("theme", "light"),
    )


def json_to_docx_executor(args: dict) -> str:
    return json_to_docx(
        json_path=args["json_path"],
        output_path=args["output_path"],
        template=args.get("template", ""),
    )


def data_to_xlsx_executor(args: dict) -> str:
    return data_to_xlsx(
        data_path=args["data_path"],
        output_path=args["output_path"],
        sheet_name=args.get("sheet_name", "Sheet1"),
    )


def markdown_to_pptx_executor(args: dict) -> str:
    return markdown_to_pptx(
        md_path=args["md_path"],
        output_path=args["output_path"],
        template=args.get("template", ""),
    )


def text_to_file_executor(args: dict) -> str:
    """Create a text file with specified content."""
    content = args.get("content", args.get("text", ""))
    output_path = args.get("output_path", args.get("path", args.get("file_path", "")))
    if not output_path:
        return "Error: missing output_path. Provide output_path in args."
    encoding = args.get("encoding", "utf-8")
    return text_to_file(content, output_path, encoding)


def html_to_file_executor(args: dict) -> str:
    """Create an HTML file with specified content."""
    html_content = args.get("html_content", args.get("content", args.get("html", "")))
    output_path = args.get("output_path", args.get("path", args.get("file_path", "")))
    if not output_path:
        return "Error: missing output_path. Provide output_path in args."
    return html_to_file(html_content, output_path)

